import os
import glob
import time
from collections import defaultdict
from typing import List, Dict
import json
import pandas as pd
import PyPDF2
import docx
from pptx import Presentation
from PIL import Image
import pytesseract
from dotenv import load_dotenv
from openai import OpenAI 


load_dotenv(override=False)

DEEPSEEK_API_KEY  = os.environ.get("DEEPSEEK_API_KEY")
DEEPSEEK_BASE_URL = os.environ.get("DEEPSEEK_BASE_URL", "https://api.deepseek.com")
DEEPSEEK_MODEL    = os.environ.get("DEEPSEEK_MODEL", "deepseek-reasoner")

PDF_FOLDER   = os.environ.get("PDF_FOLDER",   "/data/papers")
OUTPUT_EXCEL = os.environ.get("DEEPSEEK_OUTPUT", "papers-deepseek.xlsx")

if not DEEPSEEK_API_KEY:
    raise RuntimeError("DEEPSEEK_API_KEY not set. Put it in your environment or in a .env file.")

client = OpenAI(api_key=DEEPSEEK_API_KEY, base_url=DEEPSEEK_BASE_URL)


TM_COL = "TM Variant"
FIRST_COLS = [
    "Author/Year of Publication",
    "Paper Title",
    "Species Origin",
    "Production Organism",
]
DEPENDENT_COLS = [
    "Growth Factor",
    "Amino Acid Sequence",
    "Methods",
    "Experimental Conditions/Methods Description",
    "pH", "Buffer", "Protein concentration", "Tm",
]
ENGINEERED_COL = "Engineered variant?"
ALL_COLUMNS = FIRST_COLS + [TM_COL] + DEPENDENT_COLS + [ENGINEERED_COL, "PDF File Name"]



def extract_text_from_pdf(pdf_path: str) -> str:
    text = ""
    try:
        with open(pdf_path, "rb") as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
    except Exception as e:
        print(f"Error processing {pdf_path} with PyPDF2: {e}")
    try:
        from pdf2image import convert_from_path
        pages = convert_from_path(pdf_path)
        for i, page_image in enumerate(pages):
            ocr_text = pytesseract.image_to_string(page_image)
            if ocr_text.strip():
                text += f"\n[OCR extracted from image, page {i+1}]:\n{ocr_text}\n"
    except Exception as e:
        print(f"Error extracting images/OCR from {pdf_path}: {e}")
    try:
        import pdfplumber
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                tables = page.extract_tables()
                if tables:
                    for table in tables:
                        table_str = "\n[Extracted Table]:\n"
                        for row in table:
                            table_str += "\t".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
                        text += table_str + "\n"
    except Exception as e:
        print(f"Error extracting tables from {pdf_path}: {e}")
    return text

def extract_text_from_docx(docx_path: str) -> str:
    try:
        d = docx.Document(docx_path)
        return "".join(p.text + "\n" for p in d.paragraphs)
    except Exception as e:
        print(f"Error processing {docx_path}: {e}")
        return ""

def extract_text_from_pptx(pptx_path: str) -> str:
    try:
        prs = Presentation(pptx_path)
        out = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    out.append(shape.text)
        return "\n".join(out) + ("\n" if out else "")
    except Exception as e:
        print(f"Error processing {pptx_path}: {e}")
        return ""

def extract_text_from_tiff(tiff_path: str) -> str:
    try:
        image = Image.open(tiff_path)
        return pytesseract.image_to_string(image)
    except Exception as e:
        print(f"Error processing {tiff_path}: {e}")
        return ""

def extract_text_from_xlsx(xlsx_path: str) -> str:
    try:
        df = pd.read_excel(xlsx_path, engine="openpyxl")
        return df.astype(str).agg(" ".join, axis=1).str.cat(sep="\n")
    except Exception as e:
        print(f"Error processing {xlsx_path}: {e}")
        return ""

def extract_text(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    if ext == ".docx":
        return extract_text_from_docx(file_path)
    if ext == ".pptx":
        return extract_text_from_pptx(file_path)
    if ext in [".tif", ".tiff"]:
        return extract_text_from_tiff(file_path)
    if ext == ".xlsx":
        return extract_text_from_xlsx(file_path)
    print(f"Unsupported file extension: {ext} in file {file_path}")
    return ""


def call_deepseek(prompt: str, text: str, max_retries: int = 3, temperature: float = 0.1) -> str:
    truncated = text[:200_000]
    full = prompt + "\n\nTEXT:\n" + truncated
    for i in range(max_retries):
        try:
            resp = client.chat.completions.create(
                model=DEEPSEEK_MODEL,
                messages=[
                    {"role": "system", "content": "You extract scientific information precisely."},
                    {"role": "user", "content": full},
                ],
                temperature=temperature,
            )
            return resp.choices[0].message.content
        except Exception as e:
            print(f"DeepSeek attempt {i+1} failed: {e}")
            time.sleep(2 ** i)
    return ""

def build_meta_prompt(col: str) -> str:
    if col == "Production Organism":
        return (
            "Extract the production organism. Do not provide any extra information. "
            "e.g. “E.coli”, “Escherichia coli”, “P. pastoris”, “Pichia pastoris”, "
            "“S. cerevisiae”, “Saccharomyces cerevisiae”, “T. reesei”, “Trichoderma reesei”, "
            "“A. niger”, “Aspergillus niger”, “N. crassa”, “Neurospora crassa”, "
            "“bacteria”, “filamentous fungi”, “yeast”. At the end, cite where you got your information."
        )
    base = f"Extract the {col} related to this paper. Provide only the {col} as plain text."
    if col not in ("Author/Year of Publication", "Paper Title"):
        base += " Cite section or figure."
    return base

def build_variant_prompt(dep: str, variant: str) -> str:
    if dep == "Growth Factor":
        return (
            f"For variant “{variant}”, extract only the “{dep}”. Provide only the “{dep}” as one or two words. "
            "Possible choices: EGF, HGF, IGF-1, IGF-1-LR3, Insulin, Transferrin, Albumin, Shh, Wnt3A, BMP-4, "
            "FGF-2, IGF-2, IL4, IL6, TGFB1, TGFB2, TGFB3, FGF-6, Activin A, BMP-1, BMP-2, BMP-9, FGF-1, FGF-4, "
            "FGF2-STAB, FGF21, HGF NK1, PDGF-AA, PDGF-BB, TGFa, VEGF-A. Do not provide any extra information."
        )
    if dep == "Amino Acid Sequence":
        return (
            f"I am interested in the variant \"{variant}\".\n"
            "Please do the following:\n"
            "1. Extract the UniProt accession ID of the wild-type protein (e.g. P12345).\n"
            "2. List all mutations for this variant, giving the one-letter code change plus residue number, "
            "and—if mentioned—a brief location note.\n\n"
            "Output **only** plain text in this exact format:\n"
            "UniProt ID: <ID or “not found”>\n"
            "Mutations:\n"
            "- <Mutation1> (<location note if any>)\n"
            "- <Mutation2> (<location note if any>)\n\n"
            f"Variant context: {variant}\n"
        )
    if dep == "Methods":
        return (
            f"Extract the {dep}. Provide only the name of the method used for measuring thermal stability "
            f"of the variant “{variant}” as plain text. e.g. “Differential Scanning Calorimetry (DSC)”; "
            f"“Circular Dichroism (CD)”."
        )
    return (
        f"For variant “{variant}”, extract only the {dep}. Do not provide any extra information. "
        f"At the end, tell where you got your information"
    )

def build_engineered_prompt(variant: str) -> str:
    return f"For variant “{variant}”, is it engineered? Answer Yes or No. Cite location."

def extract_tm_variants(text: str) -> List[str]:
    prompt = (
        "List all distinct thermal stability (Tm) variants reported in the paper(s). "
        "Return JSON with a top-level array named 'variants'. Example:\n"
        "{ \"variants\": [\"EGF WT\", \"EGF L52I\"] }\n"
        "If none, return {\"variants\": []}. Do not invent variants."
    )
    ans = call_deepseek(prompt, text)
    # Try JSON parse first
    try:
        data = json.loads(ans)
        if isinstance(data, dict) and isinstance(data.get("variants"), list):
            vals = [str(x).strip() for x in data["variants"] if str(x).strip()]
        elif isinstance(data, list):
            vals = [str(x).strip() for x in data if str(x).strip()]
        else:
            raise ValueError("Unexpected JSON shape")
    except Exception:
        # Fallback: parse lines/bullets
        vals = [line.strip("•- \t").strip() for line in ans.splitlines() if line.strip()]
    seen, out = set(), []
    for v in vals:
        if v and v not in seen:
            seen.add(v); out.append(v)
    return out


def group_files_by_key(pdf_folder: str) -> Dict[str, List[str]]:
    groups: Dict[str, List[str]] = defaultdict(list)
    for p in glob.glob(os.path.join(pdf_folder, "*.*")):
        key = os.path.basename(p).split('-')[0].split('_')[0].strip()
        groups[key].append(p)
    return groups

def main():
    df = pd.DataFrame(columns=ALL_COLUMNS)
    file_groups = group_files_by_key(PDF_FOLDER)

    for key, files in file_groups.items():
        print(f"\nProcessing {key} ({len(files)} files)")
        combined = ""
        for f in files:
            combined += extract_text(f) + "\n"
        if not combined.strip():
            continue

        meta = {"PDF File Name": "; ".join(os.path.basename(f) for f in files)}
        for col in FIRST_COLS:
            meta[col] = call_deepseek(build_meta_prompt(col), combined)
            time.sleep(1)

        variants = extract_tm_variants(combined)
        if not variants:
            print(f"⚠️ No Tm variants detected for {key}, skipping.")
            continue
        print("variants:", variants)

        # ---- Per-variant rows ----
        for var in variants:
            row = {c: meta.get(c, "") for c in FIRST_COLS}
            row["PDF File Name"] = meta["PDF File Name"]
            row[TM_COL] = var

            for dep in DEPENDENT_COLS:
                q = build_variant_prompt(dep, var)
                row[dep] = call_deepseek(q, combined)
                time.sleep(1)

            row[ENGINEERED_COL] = call_deepseek(build_engineered_prompt(var), combined)
            time.sleep(1)

            df = pd.concat([df, pd.DataFrame([row])], ignore_index=True)

        df.to_excel(OUTPUT_EXCEL, index=False)
        print(f"Saved {len(variants)} rows for {key}")

    print("Done. Results in", OUTPUT_EXCEL)

if __name__ == "__main__":
    main()
