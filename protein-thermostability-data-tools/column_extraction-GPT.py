import os
import glob
import time
from collections import defaultdict
import openai
import PyPDF2
import pandas as pd
import docx
from pptx import Presentation
from PIL import Image
import pytesseract
from openai import error as openai_error
import json
from openai.error import RateLimitError
from dotenv import load_dotenv
from typing import Dict, List, Optional, Tuple




load_dotenv(override=False)  # reads .env if available
OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY")
if not OPENAI_API_KEY:
    raise RuntimeError("OPENAI_API_KEY not set. Put it in your environment or in a .env file.")
openai.api_key = OPENAI_API_KEY

# allow overriding paths via env
PDF_FOLDER   = os.environ.get("PDF_FOLDER",   "/data/papers")
OUTPUT_EXCEL = os.environ.get("GPT_OUTPUT", "/data/result.xlsx")
TM_CSV       = os.environ.get("TM_CSV",       "/data/manual-checked.csv")
MODEL_NAME    = os.environ.get("GPT_MODEL", "gpt-4.1")

# Columns to be extracted from papers
TM_COL = "TM Variant"
FIRST_COLS = [
    "Author/Year of Publication",
    "Paper Title",
    "Species Origin",
    "Production Organism",
]
DEPENDENT_COLS = [
    "Growth Factor",
    "Amino Acid Sequence",
    "Methods",
    "Experimental Conditions/Methods Description",
    "pH", "Buffer", "Protein concentration", "Tm",
]
ENGINEERED_COL = "Engineered variant?"
ALL_COLUMNS = FIRST_COLS + [TM_COL] + DEPENDENT_COLS + [ENGINEERED_COL, "PDF File Name"]

def system_prompt() -> Dict[str, str]:
    return {
        "role": "system",
        "content": "You are an expert in extracting scientific details from papers."
    }

def load_text_prompt(combined_text: str) -> Dict[str, str]:
    return {
        "role": "user",
        "content": f"Here is the full text of the paper(s):\n\n{combined_text}\n\nPlease remember this text."
    }

def confirm_loaded_prompt() -> Dict[str, str]:
    return {
        "role": "user",
        "content": "Confirm you have loaded the text. We want to extract columns from it and create a dataset."
    }

def meta_prompt(col: str) -> str:
    """Paper-level metadata prompt (for FIRST_COLS)."""
    if col == "Production Organism":
        return (
            "extract the production organism. Do not provide any extra information. "
            "e.g. “E.coli”, “Escherichia coli”, “P. pastoris”, “Pichia pastoris”, "
            "“S. cerevisiae”, “Saccharomyces cerevisiae”, “T. reesei”, “Trichoderma reesei”, "
            "“A. niger”, “Aspergillus niger”, “N. crassa”, “Neurospora crassa”, "
            "“bacteria”, “filamentous fungi”, “yeast”. At the end, tell where you got your information"
        )
    base = f"Extract the {col} related to this paper. Provide only the {col} as plain text."
    if col not in ("Author/Year of Publication", "Paper Title"):
        base += " Cite section or figure."
    return base

def variant_prompt(dep: str, variant: str) -> str:
    """Variant-level prompt for DEPENDENT_COLS."""
    if dep == "Growth Factor":
        return (
            f"For variant “{variant}”, extract only the “{dep}”. Provide only the “{dep}” as one or two words. "
            "Possible choices: EGF, HGF, IGF-1, IGF-1-LR3, Insulin, Transferrin, Albumin, Shh, Wnt3A, BMP-4, "
            "FGF-2, IGF-2, IL4, IL6, TGFB1, TGFB2, TGFB3, FGF-6, Activin A, BMP-1, BMP-2, BMP-9, FGF-1, FGF-4, "
            "FGF2-STAB, FGF21, HGF NK1, PDGF-AA, PDGF-BB, TGFa, VEGF-A. Do not provide any extra information."
        )
    if dep == "Amino Acid Sequence":
        return (
            f"I am interested in the variant \"{variant}\".\n"
            "Please do the following:\n"
            "1. Extract the UniProt accession ID of the wild-type protein (e.g. P12345).\n"
            "2. List all mutations for this variant, giving the one-letter code change plus residue number, "
            "and—if mentioned—a brief location note.\n\n"
            "Output **only** plain text in this exact format:\n"
            "UniProt ID: <ID or “not found”>\n"
            "Mutations:\n"
            "- <Mutation1> (<location note if any>)\n"
            "- <Mutation2> (<location note if any>)\n\n"
            f"Variant context: {variant}\n"
        )
    if dep == "Methods":
        return (
            f"Extract the {dep}. Provide only the name of the method used for measuring thermal stability "
            f"of the variant “{variant}” as plain text. e.g. “Differential Scanning Calorimetry (DSC)”; "
            f"“Circular Dichroism (CD)”."
        )
    # default
    return (
        f"For variant “{variant}”, extract only the {dep}. Do not provide any extra information. "
        f"At the end, tell where you got your information"
    )

def engineered_prompt(variant: str) -> Dict[str, str]:
    return {
        "role": "user",
        "content": f"For variant “{variant}”, is it engineered? Answer Yes or No."
    }



def extract_text_from_pdf(pdf_path):
    text = ""
    try:
        with open(pdf_path, "rb") as file:
            reader = PyPDF2.PdfReader(file)     # Extract text with PyPDF2
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
    except Exception as e:
        print(f"Error processing {pdf_path} with PyPDF2: {e}")
    try:
        from pdf2image import convert_from_path     # Extract images and perform OCR
        pages = convert_from_path(pdf_path)
        for i, page_image in enumerate(pages):
            ocr_text = pytesseract.image_to_string(page_image)
            if ocr_text.strip():
                text += f"\n[OCR extracted from image, page {i+1}]:\n{ocr_text}\n"
    except Exception as e:
        print(f"Error extracting images/OCR from {pdf_path}: {e}")
    try:
        import pdfplumber      # Extract tables with pdfplumber
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                tables = page.extract_tables()
                if tables:
                    for table in tables:
                        table_str = "\n[Extracted Table]:\n"
                        for row in table:
                            table_str += "\t".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
                        text += table_str + "\n"
    except Exception as e:
        print(f"Error extracting tables from {pdf_path}: {e}")
    return text


def extract_text_from_docx(docx_path):
    text = ""
    try:
        doc = docx.Document(docx_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        print(f"Error processing {docx_path}: {e}")
    return text


def extract_text_from_pptx(pptx_path):
    text = ""
    try:
        prs = Presentation(pptx_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error processing {pptx_path}: {e}")
    return text


def extract_text_from_tiff(tiff_path):
    text = ""
    try:
        image = Image.open(tiff_path)
        text = pytesseract.image_to_string(image)
    except Exception as e:
        print(f"Error processing {tiff_path}: {e}")
    return text


def extract_text_from_xlsx(xlsx_path):
    text = ""
    try:
        df = pd.read_excel(xlsx_path, engine="openpyxl")
        text = df.astype(str).agg(" ".join, axis=1).str.cat(sep="\n")
    except Exception as e:
        print(f"Error processing {xlsx_path}: {e}")
    return text


def extract_text(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext == ".docx":
        return extract_text_from_docx(file_path)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_path)
    elif ext in [".tif", ".tiff"]:
        return extract_text_from_tiff(file_path)
    elif ext == ".xlsx":
        return extract_text_from_xlsx(file_path)
    else:
        print(f"Unsupported file extension: {ext} in file {file_path}")
        return ""


"""
Call GPT O3-mini API to extract the columns from the text.
This method creates a chat and gives the paper as the history of the chat, but not the previous messages
"""
def ask_chat(history, retries=3):
    backoff = 1
    for i in range(retries):
        try:
            resp = openai.ChatCompletion.create(
                model=MODEL_NAME,
                messages=history
            )
            return resp.choices[0].message.content.strip()
        except RateLimitError:
            wait = 60  # seconds
            print(f"🔔 Rate limit hit—sleeping {wait}s before retry #{i+1}")
            time.sleep(wait)
        except (openai_error.APIError, openai_error.APIConnectionError) as e:
             print(f"Chat API error (attempt {i+1}): {e}")
             time.sleep(backoff)
             backoff *= 2
        except Exception as e:
            print(f"Chat attempt {i+1} failed: {e}")
            time.sleep(backoff)
            backoff *= 2
    return "[Error]"


def choose_tm(row):
    flag = row["Correct or not?"].strip().lower()
    if "correct" in flag:
        return row["TM Variant"].strip()
    if any(k in flag for k in ("revised", "added")):
        return row["Revised or Not."].strip()
    return None

def group_files_by_key(pdf_folder: str) -> Dict[str, List[str]]:
    groups: Dict[str, List[str]] = defaultdict(list)
    for path in glob.glob(os.path.join(pdf_folder, "*.*")):
        key = os.path.basename(path).split('-')[0].split('_')[0].strip()
        groups[key].append(path)
    return groups



def main():
    # Manual Tm list grouped by key
    manual_tm = pd.read_csv(TM_CSV, dtype=str).fillna("")
    manual_tm["chosen_tm"] = manual_tm.apply(choose_tm, axis=1)
    manual_tm = manual_tm.dropna(subset=["chosen_tm"])
    tm_by_key = (
        manual_tm
        .groupby("Key")["chosen_tm"]
        .apply(list)
        .to_dict()
    )

    # gets the paper and supplementary 
    # ( All the papers were named like this [author year - nameOfPaper] and [author year] was the same among the paper and supplementary)
    file_groups = group_files_by_key(PDF_FOLDER)
    master_df = pd.DataFrame(columns=ALL_COLUMNS)
    for key, files in file_groups.items():
        print(f"\n--- Processing {key} ({len(files)} files) ---")
        combined = ""
        for f in files:
            combined += extract_text(f) + "\n"
        if not combined.strip():
            continue

        # Init chat memory
        messages = [system_prompt(), load_text_prompt(combined)]
        print(ask_chat(messages + [confirm_loaded_prompt()]))
        meta = {"PDF File Name": "; ".join(os.path.basename(f) for f in files)}
        for col in FIRST_COLS:
            prompt = meta_prompt(col)
            meta[col] = ask_chat(messages + [{"role": "user", "content": prompt}])
            time.sleep(1)
        variants = tm_by_key.get(key, []).copy()
        if not variants:
            print(f"⚠️ No manual Tₘ found for {key}, skipping.")
            continue

        for var in variants:
            row = {c: meta[c] for c in FIRST_COLS}
            row["PDF File Name"] = meta["PDF File Name"]
            row[TM_COL] = var
            for dep in DEPENDENT_COLS:
                q = variant_prompt(dep, var)
                row[dep] = ask_chat(messages + [{"role": "user", "content": q}])
                time.sleep(1)
            # engineered?
            row[ENGINEERED_COL] = ask_chat(messages + [engineered_prompt(var)])
            time.sleep(1)

            master_df = pd.concat([master_df, pd.DataFrame([row])], ignore_index=True)

        master_df.to_excel(OUTPUT_EXCEL, index=False)
        print(f"→ Saved {len(variants)} rows for {key}")

    print("\nAll done — results in", OUTPUT_EXCEL)


if __name__ == "__main__":
    main()



